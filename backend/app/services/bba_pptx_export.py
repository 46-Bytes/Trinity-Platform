"""
Phase 3 – BBA PowerPoint Presentation Export Service.

Generates a branded .pptx file from `bba.presentation_slides` using python-pptx.
Programmatic template matching the existing Word report's navy/white BBA branding.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from app.models.bba import BBA

logger = logging.getLogger(__name__)

# ── Brand constants ──────────────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x36, 0x5D) if PPTX_AVAILABLE else None
WHITE = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5) if PPTX_AVAILABLE else None
DARK_GREY = RGBColor(0x33, 0x33, 0x33) if PPTX_AVAILABLE else None
ACCENT_BLUE = RGBColor(0x2B, 0x6C, 0xB0) if PPTX_AVAILABLE else None

FONT_FAMILY = "Calibri"
SLIDE_WIDTH = Inches(13.333) if PPTX_AVAILABLE else None   # 16:9
SLIDE_HEIGHT = Inches(7.5) if PPTX_AVAILABLE else None


class BBAPptxExporter:
    """
    Exports BBA presentation slides to a branded PowerPoint file.
    """

    def __init__(self):
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PowerPoint export. "
                "Install it with: pip install python-pptx"
            )

    def generate_presentation(self, bba: BBA) -> bytes:
        """
        Build a complete .pptx deck from bba.presentation_slides.

        Returns:
            Bytes of the generated PowerPoint file.
        """
        logger.info("[BBA PPTX] Generating presentation for BBA %s", bba.id)

        slides_data = (bba.presentation_slides or {}).get("slides", [])
        if not slides_data:
            raise ValueError("No presentation slides exist. Generate slides first.")

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        # Use blank layout for full control
        blank_layout = prs.slide_layouts[6]  # Blank

        for slide_data in slides_data:
            slide_type = slide_data.get("type", "")

            if slide_type == "title":
                self._add_title_slide(prs, blank_layout, slide_data)
            elif slide_type == "executive_summary":
                self._add_content_slide(prs, blank_layout, slide_data)
            elif slide_type == "structure":
                self._add_content_slide(prs, blank_layout, slide_data)
            elif slide_type == "recommendation":
                self._add_recommendation_slide(prs, blank_layout, slide_data)
            elif slide_type == "timeline":
                self._add_timeline_slide(prs, blank_layout, slide_data)
            elif slide_type == "next_steps":
                self._add_content_slide(prs, blank_layout, slide_data)
            else:
                # Generic fallback
                self._add_content_slide(prs, blank_layout, slide_data)

        # Save to bytes
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        logger.info("[BBA PPTX] Presentation generated (%d slides)", len(slides_data))
        return buf.getvalue()

    # ─── Slide builders ──────────────────────────────────────────────────

    def _add_title_slide(
        self, prs: "Presentation", layout, data: Dict[str, Any]
    ) -> None:
        """Full navy background, centred white title + subtitle."""
        slide = prs.slides.add_slide(layout)

        # Navy background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = NAVY

        # Title text
        title = data.get("title", "Diagnostic Findings & Recommendations")
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(2.2), Inches(11.333), Inches(1.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle = data.get("subtitle", "")
        if subtitle:
            txBox2 = slide.shapes.add_textbox(
                Inches(1), Inches(4.0), Inches(11.333), Inches(1.0)
            )
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.name = FONT_FAMILY
            p2.font.size = Pt(20)
            p2.font.color.rgb = WHITE
            p2.alignment = PP_ALIGN.CENTER

        # Thin accent line
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4), Inches(3.8),
            Inches(5.333), Pt(3),
        )
        line_shape = slide.shapes[-1]
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = WHITE
        line_shape.line.fill.background()

    def _add_content_slide(
        self, prs: "Presentation", layout, data: Dict[str, Any]
    ) -> None:
        """Navy title bar at top, white body with bullets."""
        slide = prs.slides.add_slide(layout)
        self._add_title_bar(slide, data.get("title", ""))

        bullets = data.get("bullets") or []
        if bullets:
            self._add_bullet_list(
                slide, bullets,
                left=Inches(0.8), top=Inches(1.7),
                width=Inches(11.733), height=Inches(5.2),
            )

    def _add_recommendation_slide(
        self, prs: "Presentation", layout, data: Dict[str, Any]
    ) -> None:
        """
        Three-column layout for recommendation slides:
        Finding | Recommendation | Outcome
        """
        slide = prs.slides.add_slide(layout)
        self._add_title_bar(slide, data.get("title", ""))

        col_width = Inches(3.78)
        col_gap = Inches(0.15)
        top = Inches(1.7)
        height = Inches(5.2)
        start_left = Inches(0.8)

        sections = [
            ("Finding", data.get("finding") or []),
            ("Recommendation", data.get("recommendation_bullets") or []),
            ("Outcome", data.get("outcome") or []),
        ]

        for i, (heading, bullets) in enumerate(sections):
            left = start_left + i * (col_width + col_gap)

            # Column heading
            hBox = slide.shapes.add_textbox(left, top, col_width, Inches(0.45))
            htf = hBox.text_frame
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hp.text = heading
            hp.font.name = FONT_FAMILY
            hp.font.size = Pt(14)
            hp.font.bold = True
            hp.font.color.rgb = NAVY

            # Bullets
            if bullets:
                self._add_bullet_list(
                    slide, bullets,
                    left=left, top=top + Inches(0.55),
                    width=col_width, height=height - Inches(0.55),
                    font_size=Pt(12),
                )

    def _add_timeline_slide(
        self, prs: "Presentation", layout, data: Dict[str, Any]
    ) -> None:
        """Table-style timeline slide."""
        slide = prs.slides.add_slide(layout)
        self._add_title_bar(slide, data.get("title", "Implementation Timeline"))

        rows_data = data.get("rows") or []
        if not rows_data:
            return

        num_rows = len(rows_data) + 1  # +1 for header
        num_cols = 4  # Rec, Title, Timing, Key Outcome

        tbl_left = Inches(0.8)
        tbl_top = Inches(1.7)
        tbl_width = Inches(11.733)
        tbl_height = Inches(0.4) * num_rows

        table_shape = slide.shapes.add_table(
            num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height
        )
        table = table_shape.table

        # Column widths
        col_widths = [Inches(0.8), Inches(4.5), Inches(2.5), Inches(3.933)]
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

        # Header row
        headers = ["#", "Recommendation", "Timing", "Key Outcome"]
        for ci, header_text in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = header_text
            self._style_table_cell(cell, bold=True, bg=NAVY, fg=WHITE, size=Pt(11))

        # Data rows
        for ri, row in enumerate(rows_data, start=1):
            vals = [
                str(row.get("rec", ri)),
                row.get("title", ""),
                row.get("timing", ""),
                row.get("outcome", ""),
            ]
            for ci, val in enumerate(vals):
                cell = table.cell(ri, ci)
                cell.text = val
                bg = LIGHT_GREY if ri % 2 == 0 else WHITE
                self._style_table_cell(cell, bold=False, bg=bg, fg=DARK_GREY, size=Pt(10))

    # ─── Shared helpers ──────────────────────────────────────────────────

    def _add_title_bar(self, slide, title_text: str) -> None:
        """Add a navy title bar across the top of the slide."""
        # Navy rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            SLIDE_WIDTH, Inches(1.2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()

        # Title text
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_FAMILY
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT

    def _add_bullet_list(
        self,
        slide,
        bullets: List[str],
        left,
        top,
        width,
        height,
        font_size=None,
    ) -> None:
        """Add a bulleted list to the slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.name = FONT_FAMILY
            p.font.size = font_size or Pt(14)
            p.font.color.rgb = DARK_GREY
            p.space_after = Pt(6)

    @staticmethod
    def _style_table_cell(cell, bold: bool, bg, fg, size) -> None:
        """Apply consistent styling to a table cell."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_FAMILY
            paragraph.font.size = size
            paragraph.font.bold = bold
            paragraph.font.color.rgb = fg

        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Reduce cell margins for tighter layout
        cell.margin_left = Pt(6)
        cell.margin_right = Pt(6)
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
