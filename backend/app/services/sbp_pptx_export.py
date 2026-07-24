"""
Strategic Business Plan PowerPoint Exporter
Generates .pptx files from generated slide content.
"""
import io
import logging
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.models.strategic_business_plan import StrategicBusinessPlan

logger = logging.getLogger(__name__)


class SBPPptxExporter:
    """Generates a .pptx presentation from the plan's slide content."""

    def generate_pptx(self, plan: StrategicBusinessPlan) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slides_data = []
        if plan.presentation_slides and isinstance(plan.presentation_slides, dict):
            slides_data = plan.presentation_slides.get("slides", [])

        for slide_data in slides_data:
            slide_type = slide_data.get("type", "section_overview")

            if slide_type == "title":
                self._add_title_slide(prs, slide_data, plan)
            else:
                self._add_content_slide(prs, slide_data)

        # Save
        buffer = io.BytesIO()
        prs.save(buffer)
        pptx_bytes = buffer.getvalue()
        logger.info(f"Generated SBP presentation for plan {plan.id} ({len(pptx_bytes)} bytes)")
        return pptx_bytes

    def _add_title_slide(self, prs: Presentation, data: Dict[str, Any], plan: StrategicBusinessPlan):
        layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(layout)

        if slide.placeholders[0]:
            slide.placeholders[0].text = data.get("title", "Strategic Business Plan")
        if len(slide.placeholders) > 1 and slide.placeholders[1]:
            subtitle = data.get("subtitle", plan.client_name or "")
            slide.placeholders[1].text = subtitle

    def _add_content_slide(self, prs: Presentation, data: Dict[str, Any]):
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        if slide.placeholders[0]:
            slide.placeholders[0].text = data.get("title", "")

        bullets = data.get("bullets") or []
        if len(slide.placeholders) > 1:
            content_ph = slide.placeholders[1]
            if bullets:
                tf = content_ph.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.paragraphs[0].text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
            else:
                sp = content_ph._element
                sp.getparent().remove(sp)


def get_sbp_pptx_exporter() -> SBPPptxExporter:
    return SBPPptxExporter()
